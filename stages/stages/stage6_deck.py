"""
Stage 6 — Deck Assembly (Enhanced Pro Design)
Builds an 8-slide PowerPoint deck using python-pptx.
- Pro design: no accent lines under titles, no flat header bars
- Company-specific content: LLM-generated recommendations & why-now per company
"""

import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from config import BENCHMARKS, OUTPUT_DIR
from utils.llm import chat, chat_json


# ── Brand colours ─────────────────────────────────────────────────────────────
C_NAVY     = RGBColor(0x0D, 0x1B, 0x2A)
C_CARD     = RGBColor(0x12, 0x26, 0x3A)
C_DEEPCARD = RGBColor(0x08, 0x12, 0x1E)
C_ACCENT   = RGBColor(0x00, 0xB4, 0xD8)
C_ACCENT2  = RGBColor(0x90, 0xE0, 0xEF)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED    = RGBColor(0x7A, 0x8C, 0x9E)
C_GREEN    = RGBColor(0x06, 0xD6, 0x7E)
C_RED      = RGBColor(0xFF, 0x4D, 0x6D)
C_AMBER    = RGBColor(0xFF, 0xC3, 0x00)
C_LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFB)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def _bg(slide, color=None):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or C_NAVY


def _rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def _text(slide, text, left, top, width, height,
          size=16, bold=False, color=C_WHITE,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.margin_left  = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top   = Inches(0)
    tf.margin_bottom= Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = "Calibri"
    return txBox


def _multiline(slide, lines, left, top, width, height,
               size=13, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
        run.font.name  = "Calibri"
    return txBox


def _tag(slide, label, left, top):
    """Small muted uppercase section tag — no underlines."""
    _text(slide, label.upper(), left, top, 6, 0.3,
          size=9, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)


def _number_circle(slide, number, left, top, sz=0.42, bg=C_ACCENT, fg=C_NAVY):
    _rect(slide, left, top, sz, sz, fill_color=bg)
    _text(slide, str(number), left, top + 0.02, sz, sz - 0.04,
          size=13, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ── LLM content generators ────────────────────────────────────────────────────

GAP_LIBRARY = [
    "No schema markup (FAQ, HowTo, Product) — AI parsers can't structure your content",
    "No FAQ-format content — misses conversational query patterns AI favours",
    "Low third-party citations — few authoritative sites mention your brand",
    "Thin 'About' and 'What we do' pages — AI can't confidently summarise your offering",
    "No author authority signals — content lacks E-E-A-T signals AI models rely on",
    "Stale blog content — hasn't been refreshed; AI models deprioritise old content",
    "No comparison/alternative content — missing the 'X vs Y' queries where buyers decide",
    "Competitor-heavy SERPs — competitors dominate the top citations, crowding you out",
    "Weak meta descriptions — not optimised for AI snippet extraction",
    "No video/multimedia — missing signals that improve multi-modal AI visibility",
    "Missing use-case landing pages — buyers can't find category-specific proof",
    "No thought leadership content — brand not cited in opinion/research queries",
]


def _select_gaps(profile, audit):
    context = (
        f"Company: {profile.get('company_name')}\n"
        f"Category: {profile.get('category')}\n"
        f"One-liner: {profile.get('one_liner')}\n"
        f"Has schema markup: {profile.get('schema_markup')}\n"
        f"Has FAQ block: {profile.get('faq_block')}\n"
        f"AI citation rate: {audit.get('citation_rate_pct')}%\n"
        f"Top competitors cited: {[c['name'] for c in audit.get('top_competitors', [])]}\n"
    )
    gaps_list = "\n".join(f"{i+1}. {g}" for i, g in enumerate(GAP_LIBRARY))
    response = chat(
        system="You select the 4 most relevant content gaps for an AEO/GEO audit. Return ONLY 4 gap numbers (e.g. '1,3,7,9'), comma-separated, no explanation.",
        user=f"{context}\n\nGap library:\n{gaps_list}\n\nWhich 4 are most relevant? Reply with 4 numbers only."
    )
    try:
        indices = [int(x.strip()) - 1 for x in response.split(",")][:4]
        return [GAP_LIBRARY[i] for i in indices if 0 <= i < len(GAP_LIBRARY)]
    except Exception:
        return GAP_LIBRARY[:4]


def _generate_recommendations(profile, gaps):
    """3 company-specific 90-day action phases via LLM."""
    gaps_text = "\n".join(f"- {g}" for g in gaps)
    try:
        result = chat_json(
            system=(
                "You are a senior AEO/GEO strategist. Generate a 3-phase 90-day action plan. "
                "Return a JSON array with exactly 3 objects, each with keys: "
                "'phase' (e.g. 'Month 1 — Foundation'), 'action' (1 bold sentence max 18 words), "
                "'detail' (1 sentence max 25 words). "
                "Make recommendations SPECIFIC to the company's category and gaps. "
                "Return only valid JSON array, no markdown."
            ),
            user=(
                f"Company: {profile.get('company_name')}\n"
                f"Category: {profile.get('category')}\n"
                f"One-liner: {profile.get('one_liner')}\n"
                f"Target customers: {profile.get('target_customers')}\n"
                f"Key differentiator: {profile.get('key_differentiator')}\n"
                f"Key gaps:\n{gaps_text}\n"
            )
        )
        if isinstance(result, list) and len(result) == 3:
            return result
    except Exception:
        pass
    return [
        {"phase": "Month 1 — Foundation",
         "action": "Add FAQ schema, HowTo markup, and author bios to all content pages",
         "detail": "Ensures AI parsers can extract and cite your brand in structured answers."},
        {"phase": "Month 2 — Content",
         "action": "Publish comparison pages; refresh top-traffic posts with AI-friendly formatting",
         "detail": "Targets high-intent decision queries where AI models select citations."},
        {"phase": "Month 3 — Authority",
         "action": "Earn 3+ third-party citations; launch thought leadership for uncited queries",
         "detail": "Builds off-site authority signals that AI citation algorithms prioritise."},
    ]


def _generate_why_now(profile, audit):
    """4 company-specific urgency bullets via LLM."""
    try:
        result = chat_json(
            system=(
                "You write urgency-focused sales copy for AEO/GEO audits. "
                "Generate exactly 4 punchy sentences explaining WHY this company needs to act NOW. "
                "Each sentence max 18 words. Be specific to their category and competitive position. "
                "Return a JSON array of 4 strings only, no markdown."
            ),
            user=(
                f"Company: {profile.get('company_name')}\n"
                f"Category: {profile.get('category')}\n"
                f"AI citation rate: {audit.get('citation_rate_pct', 0)}%\n"
                f"Top competitor cited: {audit.get('top_competitors', [{}])[0].get('name', 'competitors') if audit.get('top_competitors') else 'competitors'}\n"
                f"Differentiator: {profile.get('key_differentiator', '')}\n"
            )
        )
        if isinstance(result, list) and len(result) == 4:
            return result
    except Exception:
        pass
    return [
        "AI search habits are forming now — early movers lock in citation positions permanently",
        "Competitor brands are investing in AEO today; the leapfrog window is 6–12 months",
        "Google's AI Overviews algorithm is maturing — influence now costs less than in 2026",
        "Every month of invisibility sends buyers directly to your top competitors",
    ]


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide1_cover(prs, profile):
    sl = _slide(prs)
    _bg(sl, C_NAVY)

    # Right panel (darker block)
    _rect(sl, 7.1, 0, 2.9, 7.5, fill_color=RGBColor(0x06, 0x12, 0x1E))
    # Right cyan edge strip
    _rect(sl, 9.75, 0, 0.25, 7.5, fill_color=C_ACCENT)

    # Section tag (top)
    _tag(sl, "AEO / GEO Audit Report", 0.55, 0.32)

    # Company name - large, anchored top-left
    company = profile.get("company_name", "Your Company")
    _text(sl, company, 0.55, 0.65, 6.3, 1.25,
          size=46, bold=True, color=C_WHITE)

    # Tagline directly below (tight)
    _text(sl, "Why You're Invisible in AI Search — and How to Fix It",
          0.55, 1.95, 6.2, 0.6, size=17, color=C_ACCENT2, italic=True)

    # Category pill
    cat = profile.get("category", "")
    if cat:
        pill_w = min(len(cat) * 0.135 + 0.5, 5.0)
        _rect(sl, 0.55, 2.75, pill_w, 0.42, fill_color=C_ACCENT)
        _text(sl, cat, 0.7, 2.77, pill_w - 0.25, 0.38,
              size=13, bold=True, color=C_NAVY)

    # One-liner
    one_liner = profile.get("one_liner", "")
    if one_liner:
        _text(sl, one_liner, 0.55, 3.38, 6.2, 0.5, size=13, color=C_MUTED, italic=True)

    # Mid-section divider
    _rect(sl, 0.55, 4.05, 6.0, 0.025, fill_color=RGBColor(0x1A, 0x3A, 0x54))

    # Content preview badges
    badges = [("Gap Analysis", C_RED), ("Competitor Map", C_ACCENT), ("90-Day Roadmap", C_GREEN)]
    for i, (badge, bc) in enumerate(badges):
        bx = 0.55 + i * 2.12
        _rect(sl, bx, 4.28, 1.95, 0.68, fill_color=RGBColor(0x0A, 0x18, 0x2A))
        _rect(sl, bx, 4.28, 0.05, 0.68, fill_color=bc)
        _text(sl, badge, bx + 0.16, 4.38, 1.68, 0.46, size=12, bold=True, color=C_WHITE)

    # Lower divider
    _rect(sl, 0.55, 5.18, 6.0, 0.025, fill_color=RGBColor(0x1A, 0x3A, 0x54))

    # Prepared for
    _text(sl, "Prepared for: " + profile.get("persona", "Marketing Leadership"),
          0.55, 5.38, 5.5, 0.42, size=11, color=C_MUTED)

    # AI Audit bullet points (right side feel - on left below)
    audit_points = ["8-slide analysis", "Live competitor data", "Personalised action plan"]
    for i, pt in enumerate(audit_points):
        _text(sl, "▸  " + pt, 0.55, 5.95 + i * 0.4, 5.5, 0.38, size=11, color=RGBColor(0x2A, 0x4A, 0x64))


def _slide2_market_shift(prs):
    sl = _slide(prs)
    _bg(sl, C_NAVY)

    _tag(sl, "Market Context", 0.5, 0.28)
    _text(sl, "The AI Search Shift", 0.5, 0.55, 9, 0.7, size=32, bold=True, color=C_WHITE)

    # 3 top cards
    cards = [
        ("65%",  "Zero-click searches",   "Searches ending without a click\n(SparkToro 2024)", C_ACCENT),
        ("~30%", "AI Overview coverage",  "Of all Google searches show\nAI Overviews (BrightEdge 2024)", C_ACCENT),
        ("100M+","Daily ChatGPT users",   "Actively searching for vendor\nrecommendations daily", C_GREEN),
    ]
    for i, (val, lbl, sub, accent) in enumerate(cards):
        lx = 0.5 + i * 3.1
        _rect(sl, lx, 1.5, 2.85, 2.4, fill_color=C_CARD)
        _rect(sl, lx, 1.5, 0.06, 2.4, fill_color=accent)
        _text(sl, val,  lx + 0.2, 1.62, 2.5, 0.85, size=40, bold=True, color=accent)
        _text(sl, lbl,  lx + 0.2, 2.48, 2.5, 0.45, size=13, bold=True, color=C_WHITE)
        _multiline(sl, sub.split("\n"), lx + 0.2, 2.97, 2.55, 0.75, size=11, color=C_MUTED)

    # Wide bottom card
    _rect(sl, 0.5, 4.15, 9.2, 1.85, fill_color=C_CARD)
    _rect(sl, 0.5, 4.15, 0.06, 1.85, fill_color=C_AMBER)
    _text(sl, "80%", 0.75, 4.25, 2.0, 0.9, size=48, bold=True, color=C_AMBER)
    _text(sl, "Citation Concentration",    2.95, 4.28, 6.5, 0.5, size=18, bold=True, color=C_WHITE)
    _text(sl, "Top 3 cited brands capture 80% of all AI answer impressions — the winner-takes-most dynamic is already here.",
          2.95, 4.85, 6.6, 0.85, size=13, color=C_MUTED)


def _slide3_where_you_stand(prs, profile, audit):
    sl = _slide(prs)
    _bg(sl, C_LIGHT_BG)

    _tag(sl, "Current State", 0.5, 0.25)
    _text(sl, "Where You Stand Today", 0.5, 0.52, 9, 0.65, size=30, bold=True, color=C_NAVY)

    # Column headers
    _rect(sl, 0.5, 1.35, 4.3, 0.55, fill_color=C_NAVY)
    _text(sl, "Traditional SEO", 0.68, 1.4, 4.0, 0.45, size=14, bold=True, color=C_WHITE)

    _rect(sl, 5.2, 1.35, 4.3, 0.55, fill_color=C_ACCENT)
    _text(sl, "AI Visibility", 5.38, 1.4, 4.0, 0.45, size=14, bold=True, color=C_NAVY)

    seo_items = [
        (f"Schema markup: {'✓ Present' if profile.get('schema_markup') else '✗ Missing'}", profile.get('schema_markup')),
        (f"FAQ content: {'✓ Found' if profile.get('faq_block') else '✗ Not found'}", profile.get('faq_block')),
        ("Meta descriptions: Present", True),
        ("Category: " + profile.get("category", "N/A"), True),
    ]
    rate = audit.get("citation_rate_pct", 0)
    top_comp = (audit.get("top_competitors", [{}]) or [{}])[0].get("name", "Unknown")
    company = profile.get("company_name", "Company")
    ai_items = [
        (f"AI citation rate: {rate}% ({audit.get('cited_count',0)}/{audit.get('total_queries',0)} queries)", rate > 40),
        (f"Top cited competitor: {top_comp}", False),
        (f"Queries without {company}: {audit.get('not_cited_count', 0)}", False),
        (f"Data source: {audit.get('data_source', 'simulation')}", True),
    ]

    for i, (item, good) in enumerate(seo_items):
        ty = 2.1 + i * 0.77
        _rect(sl, 0.5, ty, 4.3, 0.64, fill_color=C_WHITE)
        col = C_GREEN if (good and "✓" in item) else C_RED if "✗" in item else C_NAVY
        _rect(sl, 0.5, ty, 0.05, 0.64, fill_color=col)
        _text(sl, item, 0.68, ty + 0.12, 3.9, 0.42, size=12, color=C_NAVY)

    for i, (item, good) in enumerate(ai_items):
        ty = 2.1 + i * 0.77
        _rect(sl, 5.2, ty, 4.3, 0.64, fill_color=C_WHITE)
        col = C_GREEN if (i == 0 and good) else C_RED if (i == 0 and not good) else C_NAVY
        _rect(sl, 5.2, ty, 0.05, 0.64, fill_color=col)
        _text(sl, item, 5.38, ty + 0.12, 3.9, 0.42, size=12, color=C_NAVY)


def _slide4_competitors(prs, audit):
    sl = _slide(prs)
    _bg(sl, C_NAVY)

    _tag(sl, "Competitive Landscape", 0.5, 0.25)
    _text(sl, "Who's Getting Cited Instead of You", 0.5, 0.52, 9.5, 0.7,
          size=30, bold=True, color=C_WHITE)

    competitors = audit.get("top_competitors", [])
    total = audit.get("total_queries", 1) or 1

    if not competitors:
        _text(sl, "No competitor citation data available.", 0.6, 2.5, 8, 0.6, size=16, color=C_MUTED)
        return

    for i, comp in enumerate(competitors[:3]):
        lx = 0.5 + i * 3.1
        pct = round(comp["mentions"] / total * 100)

        _rect(sl, lx, 1.4, 2.85, 3.75, fill_color=C_CARD)
        _rect(sl, lx, 1.4, 0.06, 3.75, fill_color=C_ACCENT)

        _text(sl, comp["name"], lx + 0.2, 1.55, 2.5, 0.55, size=17, bold=True, color=C_ACCENT)
        _text(sl, f"{pct}%", lx + 0.2, 2.18, 2.5, 0.9, size=44, bold=True, color=C_WHITE)
        _text(sl, "of AI queries", lx + 0.2, 3.1, 2.5, 0.38, size=11, color=C_MUTED)

        # Mini progress bar
        bar_total_w = 2.4
        _rect(sl, lx + 0.2, 3.6, bar_total_w, 0.1, fill_color=C_DEEPCARD)
        _rect(sl, lx + 0.2, 3.6, max(0.1, bar_total_w * pct / 100), 0.1, fill_color=C_ACCENT)

        _text(sl, f"{comp['mentions']} / {total} queries", lx + 0.2, 3.82, 2.5, 0.35, size=10, color=C_MUTED)

    cited_queries = [r for r in audit.get("results", []) if r.get("competitors_cited")]
    if cited_queries:
        q = cited_queries[0]
        _rect(sl, 0.5, 5.42, 9.2, 1.55, fill_color=C_DEEPCARD)
        _rect(sl, 0.5, 5.42, 0.06, 1.55, fill_color=C_AMBER)
        _text(sl, "REAL QUERY EXAMPLE", 0.7, 5.48, 5, 0.3, size=9, bold=True, color=C_AMBER)
        _text(sl, f"Query: \"{q['query']}\"", 0.7, 5.82, 8.7, 0.42, size=12, bold=True, color=C_WHITE)
        _text(sl, f"AI cited: {', '.join(q['competitors_cited'][:3])}", 0.7, 6.28, 8.7, 0.5, size=12, color=C_MUTED)


def _slide5_gaps(prs, gaps):
    sl = _slide(prs)
    _bg(sl, C_LIGHT_BG)

    _tag(sl, "Gap Analysis", 0.5, 0.25)
    _text(sl, "The Gaps Holding You Back", 0.5, 0.52, 9, 0.65, size=30, bold=True, color=C_NAVY)

    gap_accents = [C_RED, C_AMBER, C_RED, C_AMBER]

    for i, gap in enumerate(gaps[:4]):
        ty = 1.35 + i * 1.5
        title_part = gap.split("—")[0].strip() if "—" in gap else gap
        desc_part  = gap.split("—")[1].strip() if "—" in gap else ""

        _rect(sl, 0.5, ty, 9.2, 1.3, fill_color=C_WHITE)
        _rect(sl, 0.5, ty, 0.07, 1.3, fill_color=gap_accents[i])

        _number_circle(sl, i + 1, 0.73, ty + 0.42, sz=0.42,
                       bg=gap_accents[i], fg=C_WHITE)

        _text(sl, title_part, 1.3, ty + 0.1, 8.1, 0.45, size=14, bold=True, color=C_NAVY)
        if desc_part:
            _text(sl, desc_part, 1.3, ty + 0.6, 8.1, 0.6, size=12, color=C_MUTED)


def _slide6_recommendations(prs, profile, recs):
    sl = _slide(prs)
    _bg(sl, C_NAVY)

    _tag(sl, "Action Plan", 0.5, 0.25)
    _text(sl, "What We'd Do: 90-Day Fix", 0.5, 0.52, 9, 0.65, size=30, bold=True, color=C_WHITE)
    _text(sl, f"Tailored for {profile.get('company_name', 'your company')}",
          0.5, 1.22, 8, 0.38, size=13, color=C_MUTED, italic=True)

    phase_colors = [C_ACCENT, C_GREEN, C_AMBER]

    for i, rec in enumerate(recs[:3]):
        ty = 1.78 + i * 1.82
        _rect(sl, 0.5, ty, 9.2, 1.62, fill_color=C_CARD)
        _rect(sl, 0.5, ty, 0.07, 1.62, fill_color=phase_colors[i])

        _number_circle(sl, i + 1, 0.73, ty + 0.58, sz=0.42,
                       bg=phase_colors[i], fg=C_NAVY)

        phase_label = rec.get("phase", f"Month {i+1}")
        action_text = rec.get("action", "")
        detail_text = rec.get("detail", "")

        _text(sl, phase_label, 1.32, ty + 0.1, 7.9, 0.4,
              size=11, bold=True, color=phase_colors[i])
        _text(sl, action_text, 1.32, ty + 0.55, 7.9, 0.52, size=14, bold=True, color=C_WHITE)
        _text(sl, detail_text, 1.32, ty + 1.12, 7.9, 0.42, size=11, color=C_MUTED)


def _slide7_why_now(prs, why_now_points):
    sl = _slide(prs)
    _bg(sl, C_NAVY)

    _tag(sl, "Urgency", 0.5, 0.25)
    _text(sl, "Why Now?", 0.5, 0.52, 9, 0.65, size=34, bold=True, color=C_WHITE)

    row_bg = [C_CARD, C_DEEPCARD, C_CARD, C_DEEPCARD]

    for i, pt in enumerate(why_now_points[:4]):
        ty = 1.38 + i * 1.52
        _rect(sl, 0.5, ty, 9.2, 1.32, fill_color=row_bg[i])
        _number_circle(sl, i + 1, 0.68, ty + 0.44, sz=0.42, bg=C_ACCENT, fg=C_NAVY)
        _text(sl, pt, 1.28, ty + 0.3, 8.1, 0.72, size=14, color=C_WHITE)


def _slide8_cta(prs, profile):
    sl = _slide(prs)
    _bg(sl, C_NAVY)

    _rect(sl, 0, 0, 10, 0.07, fill_color=C_ACCENT)

    _tag(sl, "Next Step", 0.5, 0.28)

    _text(sl, "Let's build your AI visibility roadmap.",
          0.6, 0.85, 8.8, 1.5, size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    _text(sl,
          "We'll run a full audit, map your gap-to-query matrix, and deliver\n"
          "a prioritised 90-day action plan in one working session.",
          0.6, 2.75, 8.8, 1.0, size=16, color=C_MUTED, align=PP_ALIGN.CENTER)

    _rect(sl, 3.0, 3.98, 4.2, 0.75, fill_color=C_ACCENT)
    _text(sl, "Book a 30-Minute Strategy Call  →", 3.05, 4.06, 4.1, 0.6,
          size=14, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    company = profile.get("company_name", "")
    _text(sl, f"Prepared exclusively for {company}",
          0.6, 6.78, 8.8, 0.42, size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

    _rect(sl, 0, 7.43, 10, 0.07, fill_color=C_ACCENT)


# ── Main runner ───────────────────────────────────────────────────────────────

def run(profile: dict, keywords: list, audit: dict) -> str:
    print("  [Stage 6] Selecting gaps via LLM...")
    gaps = _select_gaps(profile, audit)

    profile["top_competitors_preview"] = audit.get("top_competitors", [])[:3]

    print("  [Stage 6] Generating company-specific recommendations via LLM...")
    recs = _generate_recommendations(profile, gaps)

    print("  [Stage 6] Generating company-specific why-now points via LLM...")
    why_now = _generate_why_now(profile, audit)

    print("  [Stage 6] Assembling PowerPoint deck...")
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    _slide1_cover(prs, profile)
    _slide2_market_shift(prs)
    _slide3_where_you_stand(prs, profile, audit)
    _slide4_competitors(prs, audit)
    _slide5_gaps(prs, gaps)
    _slide6_recommendations(prs, profile, recs)
    _slide7_why_now(prs, why_now)
    _slide8_cta(prs, profile)

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    safe_name = profile.get("company_name", "prospect").replace(" ", "_").lower()
    path = os.path.join(OUTPUT_DIR, f"{safe_name}_aeo_audit.pptx")
    prs.save(path)
    print(f"  [Stage 6] ✓ Deck saved: {path}")
    return path
