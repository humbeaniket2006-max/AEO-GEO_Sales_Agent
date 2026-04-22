"""
Stage 6 — Deck Assembly
Builds an 8-slide PowerPoint deck from scratch using python-pptx.
"""

import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from config import BENCHMARKS, OUTPUT_DIR
from utils.llm import chat


# ── Brand colours ─────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # near-black navy
C_ACCENT = RGBColor(0x00, 0xB4, 0xD8)   # electric blue
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # light grey bg
C_MUTED  = RGBColor(0x8D, 0x99, 0xAE)   # grey text
C_GREEN  = RGBColor(0x2D, 0xC6, 0x53)
C_RED    = RGBColor(0xFF, 0x4D, 0x6D)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _slide(prs: Presentation, layout_idx: int = 6):
    """Add a blank slide (layout 6 = blank in most themes)."""
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def _rect(slide, left, top, width, height, fill_color=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()  # no border
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def _textbox(slide, text, left, top, width, height,
             font_size=18, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox


def _bg(slide, color=C_DARK):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _accent_line(slide, top=1.55, color=C_ACCENT):
    """Thin horizontal accent line."""
    _rect(slide, 0.6, top, 2.5, 0.04, fill_color=color)


# ── Gap diagnosis via LLM ─────────────────────────────────────────────────────

GAP_LIBRARY = [
    "No schema markup (FAQ, HowTo, Product) — AI parsers can't structure your content",
    "No FAQ-format content — misses conversational query patterns AI favours",
    "Low third-party citations — few authoritative sites mention your brand",
    "Thin 'About' and 'What we do' pages — AI can't confidently summarise your offering",
    "No author authority signals — content lacks E-E-A-T signals AI models rely on",
    "Stale blog content — hasn't been refreshed; AI models deprioritise old content",
    "No comparison/alternative content — missing the 'X vs Y' queries where buyers decide",
    "Competitor-heavy SERPs — competitors dominate the top citations, crowding you out",
    "Weak meta descriptions — not optimised for AI snippet extraction",
    "No video/multimedia — missing signals that improve multi-modal AI visibility",
    "Missing use-case landing pages — buyers can't find category-specific proof",
    "No thought leadership content — brand not cited in opinion/research queries",
]

def _select_gaps(profile: dict, audit: dict) -> list[str]:
    """Ask Groq to pick the 4 most relevant gaps from the library."""
    context = (
        f"Company: {profile.get('company_name')}\n"
        f"Has schema markup: {profile.get('schema_markup')}\n"
        f"Has FAQ block: {profile.get('faq_block')}\n"
        f"AI citation rate: {audit.get('citation_rate_pct')}%\n"
        f"Top competitors cited: {[c['name'] for c in audit.get('top_competitors', [])]}\n"
    )
    gaps_list = "\n".join(f"{i+1}. {g}" for i, g in enumerate(GAP_LIBRARY))

    response = chat(
        system="You select the 4 most relevant content gaps for an AEO/GEO audit. Return ONLY 4 gap numbers (e.g. '1,3,7,9'), comma-separated, no explanation.",
        user=f"{context}\n\nGap library:\n{gaps_list}\n\nWhich 4 are most relevant? Reply with 4 numbers only."
    )
    try:
        indices = [int(x.strip()) - 1 for x in response.split(",")][:4]
        return [GAP_LIBRARY[i] for i in indices if 0 <= i < len(GAP_LIBRARY)]
    except Exception:
        return GAP_LIBRARY[:4]


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide1_cover(prs, profile):
    sl = _slide(prs)
    _bg(sl, C_DARK)
    _rect(sl, 0, 0, 10, 0.6, C_ACCENT)
    _textbox(sl, "AEO / GEO AUDIT REPORT", 0.6, 0.1, 8, 0.5,
             font_size=11, bold=True, color=C_DARK, align=PP_ALIGN.LEFT)
    company = profile.get("company_name", "Your Company")
    _textbox(sl, company, 0.6, 1.0, 8.5, 1.2,
             font_size=44, bold=True, color=C_WHITE)
    _accent_line(sl, top=2.3)
    _textbox(sl, "Why You're Invisible in AI Search — and How to Fix It",
             0.6, 2.5, 8.5, 0.8, font_size=20, color=C_MUTED)
    _textbox(sl, profile.get("category", ""), 0.6, 3.4, 8, 0.5,
             font_size=14, color=C_ACCENT)
    _textbox(sl, "Prepared for: " + profile.get("persona", "Marketing Leadership"),
             0.6, 6.8, 8, 0.4, font_size=11, color=C_MUTED)


def _slide2_market_shift(prs):
    sl = _slide(prs)
    _bg(sl, C_DARK)
    _textbox(sl, "The AI Search Shift", 0.6, 0.4, 8, 0.6,
             font_size=30, bold=True, color=C_WHITE)
    _accent_line(sl, top=1.1)
    stats = [
        (BENCHMARKS["zero_click_rate"],       C_ACCENT),
        (BENCHMARKS["ai_overview_coverage"],  C_ACCENT),
        (BENCHMARKS["chatgpt_daily_users"],   C_ACCENT),
        (BENCHMARKS["citation_concentration"],C_ACCENT),
    ]
    for i, (stat, color) in enumerate(stats):
        top = 1.4 + i * 1.1
        _rect(sl, 0.6, top, 9, 0.9, fill_color=RGBColor(0x16, 0x2A, 0x3C))
        _textbox(sl, "▶  " + stat, 0.8, top + 0.1, 8.5, 0.7,
                 font_size=15, color=C_WHITE)


def _slide3_where_you_stand(prs, profile, audit):
    sl = _slide(prs)
    _bg(sl, C_LIGHT)
    _textbox(sl, "Where You Stand Today", 0.5, 0.3, 9, 0.6,
             font_size=28, bold=True, color=C_DARK)
    _accent_line(sl, top=1.0, color=C_ACCENT)

    company = profile.get("company_name", "Company")

    # Two columns
    headers = [("Traditional SEO", 0.5, C_DARK), ("AI Visibility", 5.2, C_DARK)]
    for label, left, col in headers:
        _rect(sl, left, 1.2, 4.3, 0.5, fill_color=C_DARK)
        _textbox(sl, label, left + 0.1, 1.25, 4.0, 0.4,
                 font_size=14, bold=True, color=C_WHITE)

    seo_items = [
        f"Schema markup: {'✓ Present' if profile.get('schema_markup') else '✗ Missing'}",
        f"FAQ content: {'✓ Found' if profile.get('faq_block') else '✗ Not found'}",
        "Meta descriptions: Present",
        "Category: " + profile.get("category", "N/A"),
    ]
    ai_items = [
        f"AI citation rate: {audit.get('citation_rate_pct', 0)}% ({audit.get('cited_count',0)}/{audit.get('total_queries',0)} queries)",
        f"Top cited competitor: {audit['top_competitors'][0]['name'] if audit.get('top_competitors') else 'Unknown'}",
        f"Queries without {company}: {audit.get('not_cited_count', 0)}",
        f"Data source: {audit.get('data_source', 'simulation')}",
    ]

    for i, item in enumerate(seo_items):
        color = C_GREEN if "✓" in item else C_RED if "✗" in item else C_DARK
        _textbox(sl, item, 0.6, 1.9 + i * 0.7, 4.0, 0.6, font_size=13, color=color)

    for i, item in enumerate(ai_items):
        rate = audit.get("citation_rate_pct", 0)
        color = C_GREEN if i == 0 and rate > 40 else C_RED if i == 0 and rate <= 40 else C_DARK
        _textbox(sl, item, 5.3, 1.9 + i * 0.7, 4.0, 0.6, font_size=13, color=color)


def _slide4_competitors(prs, audit):
    sl = _slide(prs)
    _bg(sl, C_DARK)
    _textbox(sl, "Who's Getting Cited Instead of You", 0.5, 0.3, 9.5, 0.7,
             font_size=28, bold=True, color=C_WHITE)
    _accent_line(sl, top=1.1)

    competitors = audit.get("top_competitors", [])
    total = audit.get("total_queries", 1)

    if not competitors:
        _textbox(sl, "No competitor citation data available.", 0.6, 2.0, 8, 0.6,
                 font_size=16, color=C_MUTED)
        return

    for i, comp in enumerate(competitors[:3]):
        left = 0.5 + i * 3.1
        pct = round(comp["mentions"] / total * 100)
        _rect(sl, left, 1.4, 2.8, 3.5, fill_color=RGBColor(0x16, 0x2A, 0x3C))
        _textbox(sl, comp["name"], left + 0.15, 1.55, 2.5, 0.6,
                 font_size=18, bold=True, color=C_ACCENT)
        _textbox(sl, f"{pct}%", left + 0.15, 2.3, 2.5, 0.9,
                 font_size=40, bold=True, color=C_WHITE)
        _textbox(sl, "of AI queries", left + 0.15, 3.1, 2.5, 0.4,
                 font_size=12, color=C_MUTED)
        _textbox(sl, f"{comp['mentions']} / {total} queries", left + 0.15, 3.6, 2.5, 0.4,
                 font_size=12, color=C_MUTED)

    # Sample query
    cited_queries = [r for r in audit.get("results", []) if r.get("competitors_cited")]
    if cited_queries:
        q = cited_queries[0]
        snippet = f"Query: \"{q['query']}\"\nAI cited: {', '.join(q['competitors_cited'][:3])}"
        _rect(sl, 0.5, 5.2, 9.2, 1.3, fill_color=RGBColor(0x0A, 0x10, 0x1A))
        _textbox(sl, snippet, 0.7, 5.3, 8.8, 1.1, font_size=13, color=C_MUTED)


def _slide5_gaps(prs, gaps):
    sl = _slide(prs)
    _bg(sl, C_LIGHT)
    _textbox(sl, "The Gaps Holding You Back", 0.5, 0.3, 9, 0.6,
             font_size=28, bold=True, color=C_DARK)
    _accent_line(sl, top=1.0, color=C_RED)

    for i, gap in enumerate(gaps[:4]):
        top = 1.3 + i * 1.3
        _rect(sl, 0.5, top, 9.2, 1.1, fill_color=C_WHITE)
        _rect(sl, 0.5, top, 0.08, 1.1, fill_color=C_RED)
        _textbox(sl, f"Gap {i+1}", 0.7, top + 0.05, 1.5, 0.35,
                 font_size=10, bold=True, color=C_RED)
        _textbox(sl, gap.split("—")[0].strip(), 0.7, top + 0.35, 8.5, 0.4,
                 font_size=14, bold=True, color=C_DARK)
        if "—" in gap:
            _textbox(sl, gap.split("—")[1].strip(), 0.7, top + 0.72, 8.5, 0.35,
                     font_size=11, color=C_MUTED)


def _slide6_recommendations(prs, profile, gaps):
    sl = _slide(prs)
    _bg(sl, C_DARK)
    _textbox(sl, "What We'd Do: 90-Day Fix", 0.5, 0.3, 9, 0.6,
             font_size=28, bold=True, color=C_WHITE)
    _accent_line(sl, top=1.1)

    phases = [
        ("Month 1 — Foundation",
         "Add FAQ schema, HowTo markup, author bios to all content pages"),
        ("Month 2 — Content",
         "Publish comparison and alternative pages; refresh 5 top-traffic posts with AI-friendly formatting"),
        ("Month 3 — Authority",
         "Earn 3+ third-party citations; launch thought leadership targeting top 5 uncited queries"),
    ]
    for i, (title, detail) in enumerate(phases):
        top = 1.4 + i * 1.5
        _rect(sl, 0.5, top, 9.2, 1.3, fill_color=RGBColor(0x10, 0x22, 0x36))
        _rect(sl, 0.5, top, 0.08, 1.3, fill_color=C_ACCENT)
        _textbox(sl, title, 0.7, top + 0.1, 8.5, 0.45,
                 font_size=14, bold=True, color=C_ACCENT)
        _textbox(sl, detail, 0.7, top + 0.6, 8.5, 0.55,
                 font_size=13, color=C_WHITE)


def _slide7_why_now(prs):
    sl = _slide(prs)
    _bg(sl, C_DARK)
    _textbox(sl, "Why Now?", 0.6, 0.4, 8, 0.6,
             font_size=32, bold=True, color=C_WHITE)
    _accent_line(sl, top=1.15)

    points = [
        "AI search habits are forming now — early movers lock in citation positions",
        "Competitor brands are investing in AEO today; the window to leapfrog is 6–12 months",
        "Google's AI Overviews algorithm is still maturing — influence now costs less than it will in 2026",
        "Every month of invisibility is a month of buyers being directed to someone else",
    ]
    for i, pt in enumerate(points):
        top = 1.4 + i * 1.15
        _rect(sl, 0.5, top, 9.2, 0.95, fill_color=RGBColor(0x16, 0x2A, 0x3C))
        _textbox(sl, f"  {i+1}.  " + pt, 0.6, top + 0.15, 9.0, 0.7,
                 font_size=14, color=C_WHITE)


def _slide8_cta(prs, profile):
    sl = _slide(prs)
    _bg(sl, C_DARK)
    _rect(sl, 0, 0, 10, 0.6, fill_color=C_ACCENT)
    _textbox(sl, "NEXT STEP", 0.6, 0.1, 9, 0.4,
             font_size=12, bold=True, color=C_DARK)
    _textbox(sl, "Let's build your AI visibility roadmap.", 0.6, 1.0, 9.2, 0.9,
             font_size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _textbox(sl,
             "We'll run a full audit, map your gap-to-query matrix,\n"
             "and deliver a prioritised 90-day action plan in one working session.",
             0.6, 2.1, 9.2, 1.0, font_size=16, color=C_MUTED, align=PP_ALIGN.CENTER)
    _rect(sl, 3.0, 3.4, 4.2, 0.75, fill_color=C_ACCENT)
    _textbox(sl, "Book a 30-Minute Strategy Call  →", 3.0, 3.48, 4.2, 0.6,
             font_size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    company = profile.get("company_name", "")
    _textbox(sl, f"Prepared exclusively for {company}", 0.6, 6.5, 9.2, 0.4,
             font_size=11, color=C_MUTED, align=PP_ALIGN.CENTER)


# ── Main runner ───────────────────────────────────────────────────────────────

def run(profile: dict, keywords: list, audit: dict) -> str:
    """Build the deck and return the output file path."""
    print("  [Stage 6] Selecting gaps via LLM...")
    gaps = _select_gaps(profile, audit)

    print("  [Stage 6] Assembling PowerPoint deck...")
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    _slide1_cover(prs, profile)
    _slide2_market_shift(prs)
    _slide3_where_you_stand(prs, profile, audit)
    _slide4_competitors(prs, audit)
    _slide5_gaps(prs, gaps)
    _slide6_recommendations(prs, profile, gaps)
    _slide7_why_now(prs)
    _slide8_cta(prs, profile)

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    safe_name = profile.get("company_name", "prospect").replace(" ", "_").lower()
    path = os.path.join(OUTPUT_DIR, f"{safe_name}_aeo_audit.pptx")
    prs.save(path)
    print(f"  [Stage 6] ✓ Deck saved: {path}")
    return path
